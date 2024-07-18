# import streamlit as st
# from pptx import Presentation
# from dotenv import load_dotenv
# import os
# import google.generativeai as genai
# import replicate

# # Load environment variables
# load_dotenv()
# genai.configure(api_key=os.getenv("GOOGLE_API_KEY"))

# # Function to get response from Gemini Model
# def get_gemini_response(input_text, prompt):
#     model = genai.GenerativeModel('gemini-pro')
#     summary = model.generate_content([input_text, prompt])
#     return summary.text

# # Function to generate images using Replicate API
# def generate_image(prompt):

#     input = {
#         "prompt": prompt,
#         "aspect_ratio": "3:2"
#     }

#     output = replicate.run(
#         "stability-ai/stable-diffusion-3",
#         input=input
#     )
#     print("\n\n")
#     print(output)
#     print("\n\n")
#     return output

# # Streamlit App
# def main():
#     st.title("Generate summaries and images for your PPTs")

#     with st.sidebar:
#         st.subheader("Upload your .pptx files here:")
#         uploaded_file = st.file_uploader("", type="pptx")

#         slide_wise_summary=[]
#         if uploaded_file is not None:
#             st.success("File uploaded successfully!")
#             all_text = extract_text(uploaded_file)
#             summaries = get_slide_summaries(all_text)
#             slide_wise_summary=summaries
#             st.subheader("Summaries:")
#             for i, summary in enumerate(slide_wise_summary):
#                 st.write(f"Slide {i + 1} Summary: \n{summary}")

#     if uploaded_file is not None:
#         st.subheader("Images generated for each slide:")
#         for i, summary in enumerate(summaries):
#             image_url = generate_image(summary)
#             st.image(image_url, caption=f"Slide {i + 1}")

# # Function to extract text from PPTX
# def extract_text(uploaded_file):
#     prs = Presentation(uploaded_file)
#     slide_number = 1

#     all_text = []
#     for slide in prs.slides:
#         slide_text = f"Slide {slide_number}\n"
#         slide_number += 1
#         for shape in slide.shapes:
#             if hasattr(shape, "text"):
#                 slide_text += shape.text + "\n"
#         all_text.append(slide_text)

#     return all_text


# input_prompt1 = """
# You are a very descriptive and usefull text to prompt generator
# Summarize each the input text such that it is an input prompt for an image generator model.
# Understand the core concept of the input text given and respond with a prompt that is the closest match with the topic
# """


# # Function to get slide summaries
# def get_slide_summaries(all_text):
#     summaries = []
#     for slide_text in all_text:
#         summary = get_gemini_response(slide_text, input_prompt1)
#         summaries.append(summary)
#     return summaries

# if __name__ == "__main__":
#     main()


import streamlit as st
from pptx import Presentation
from dotenv import load_dotenv
import os
import google.generativeai as genai
import replicate
from PIL import Image
import io
import requests

# Load environment variables
load_dotenv()
genai.configure(api_key=os.getenv("GOOGLE_API_KEY"))

# Function to get response from Gemini Model
def get_gemini_response(input_text, prompt):
    model = genai.GenerativeModel('gemini-pro')
    summary = model.generate_content([input_text, prompt])
    return summary.text

# Function to generate images using Replicate API
def generate_image(prompt):
    input = {
        "prompt": prompt,
        "aspect_ratio": "3:2"
    }
    output = replicate.run("stability-ai/stable-diffusion-3", input=input)
    return output

# Function to convert images to JPEG format
def convert_image_to_supported_format(image_content):
    image = Image.open(io.BytesIO(image_content))
    img_byte_arr = io.BytesIO()
    image = image.convert("RGB")  # Convert to RGB to ensure compatibility with JPEG format
    image.save(img_byte_arr, format='JPEG')  # Convert to JPEG format
    img_byte_arr.seek(0)
    return img_byte_arr

# Function to replace images in the PPTX
def replace_images_in_ppt(pptx_file, image_urls, output_path):
    prs = Presentation(pptx_file)

    for slide, image_url in zip(prs.slides, image_urls):
        # st.write(image_url)
        # st.write(image_url[0])
        response = requests.get(image_url[0])
        if response.status_code == 200:
            new_img = convert_image_to_supported_format(response.content)
            
            for shape in slide.shapes:
                if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                    # Get original image properties
                    left = shape.left
                    top = shape.top
                    width = shape.width
                    height = shape.height

                    # Remove the original image
                    slide.shapes._spTree.remove(shape._element)

                    # Add the new image
                    slide.shapes.add_picture(new_img, left, top, width=width, height=height)

    prs.save(output_path)

# Streamlit App
def main():
    st.title("Generate summaries and images for your PPTs")

    with st.sidebar:
        st.subheader("Upload your .pptx files here:")
        uploaded_file = st.file_uploader("", type="pptx")

        slide_wise_summary = []
        if uploaded_file is not None:
            st.success("File uploaded successfully!")
            all_text = extract_text(uploaded_file)
            summaries = get_slide_summaries(all_text)
            slide_wise_summary = summaries
            st.subheader("Summaries:")
            for i, summary in enumerate(slide_wise_summary):
                st.write(f"Slide {i + 1} Summary: \n{summary}")

    if uploaded_file is not None:
        st.subheader("Generating images for each slide...")
        image_urls = []
        for i, summary in enumerate(slide_wise_summary):
            image_url = generate_image(summary)
            st.image(image_url, caption=f"Slide {i + 1}")
            image_urls.append(image_url)

        output_path = "C:/Users/Harish/Downloads/updated_ppt.pptx"
        replace_images_in_ppt(uploaded_file, image_urls, output_path)
        
        st.success("Images replaced successfully!")
        with open(output_path, "rb") as f:
            st.download_button(
                label="Download updated PowerPoint",
                data=f,
                file_name="updated_presentation.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )

# Function to extract text from PPTX
def extract_text(uploaded_file):
    prs = Presentation(uploaded_file)
    slide_number = 1

    all_text = []
    for slide in prs.slides:
        slide_text = f"Slide {slide_number}\n"
        slide_number += 1
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text += shape.text + "\n"
        all_text.append(slide_text)

    return all_text

input_prompt1 = """
You are a very descriptive and useful text-to-prompt generator.
Summarize each the input text such that it is an input prompt for an image generator model.
Understand the core concept of the input text given and respond with a prompt that is the closest match with the topic.
You may use your knowledge to create a promt which will result in most realistic and accurate
"""

# Function to get slide summaries
def get_slide_summaries(all_text):
    summaries = []
    for slide_text in all_text:
        summary = get_gemini_response(slide_text, input_prompt1)
        summaries.append(summary)
    return summaries

if __name__ == "__main__":
    main()
